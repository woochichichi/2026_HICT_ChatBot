# -*- coding: utf-8 -*-
"""
사내 공유용 발표자료 생성 스크립트
- 소스: README.md, docs/adr/0001~0009, docs/TROUBLESHOOTING.md
- 구성 원칙: 슬라이드당 글 최소화, '왜(Why)'가 드러나는 옵션 비교 카드 구조
- 실행: .venv/Scripts/python.exe docs/presentation/make_ppt.py
- 산출물: docs/presentation/AI코치_PoC_공유발표.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 디자인 토큰 ──────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)      # 메인 (증권사 톤)
ORANGE = RGBColor(0xE8, 0x59, 0x0C)    # 강조 (선택/결론)
TEAL = RGBColor(0x0F, 0x76, 0x6E)      # 보조 (폐쇄망/실도입)
INK = RGBColor(0x1F, 0x29, 0x37)       # 본문
GRAY = RGBColor(0x6B, 0x72, 0x80)      # 부가 설명
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)     # 카드 배경
GOOD = RGBColor(0x15, 0x80, 0x3D)      # 장점
BAD = RGBColor(0xB9, 0x1C, 0x1C)       # 단점
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ────────────────────────────────────────────────────
def _ea(run):
    """한글 폰트(East Asian typeface) 강제 지정 — font.name만으로는 라틴 폰트만 바뀜"""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def put_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines: [(text, size, bold, color, align), ...]"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name, f.size, f.bold = FONT, Pt(size), bold
        f.color.rgb = color
        _ea(run)
    return box


def rect(slide, x, y, w, h, fill, line=None, round_=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if round_:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def shape_text(shape, lines, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name, f.size, f.bold = FONT, Pt(size), bold
        f.color.rgb = color
        _ea(run)


def header(slide, kicker, title):
    """상단 공통 헤더: kicker(작은 분류 라벨) + 제목 + 포인트 라인"""
    rect(slide, Inches(0.6), Inches(0.42), Inches(0.12), Inches(0.95), ORANGE, round_=False)
    put_text(slide, Inches(0.9), Inches(0.35), Inches(11.8), Inches(0.4),
             [(kicker, 13, True, ORANGE, PP_ALIGN.LEFT)])
    put_text(slide, Inches(0.9), Inches(0.68), Inches(11.8), Inches(0.75),
             [(title, 28, True, NAVY, PP_ALIGN.LEFT)])


def takeaway(slide, text, color=NAVY, y=Inches(6.55)):
    """하단 한 줄 결론 바"""
    bar = rect(slide, Inches(0.6), y, Inches(12.13), Inches(0.62), color)
    shape_text(bar, [(text, 15, True, WHITE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)


def option_cards(slide, options, y=Inches(2.3), h=Inches(3.9)):
    """A/B/C 옵션 비교 카드. options: [(이름, 한줄설명, 장점, 단점, 선택여부), ...]"""
    n = len(options)
    gap = Inches(0.3)
    total = Inches(12.13)
    w = Emu(int((total - gap * (n - 1)) / n))
    x = Inches(0.6)
    for name, desc, pro, con, chosen in options:
        card = rect(slide, x, y, w, h, WHITE if chosen else LIGHT,
                    line=ORANGE if chosen else None)
        lines = [
            (("✅ " if chosen else "") + name, 16, True, ORANGE if chosen else INK, PP_ALIGN.LEFT),
            (desc, 12, False, GRAY, PP_ALIGN.LEFT),
            ("👍 " + pro, 12.5, False, GOOD, PP_ALIGN.LEFT),
            ("👎 " + con, 12.5, False, BAD, PP_ALIGN.LEFT),
        ]
        shape_text(card, lines)
        x = Emu(int(x + w + gap))


def flow(slide, steps, y=Inches(3.0), h=Inches(1.1), color=NAVY):
    """좌→우 화살표 플로우"""
    n = len(steps)
    gap = Inches(0.12)
    total = Inches(12.13)
    w = Emu(int((total - gap * (n - 1)) / n))
    x = Inches(0.6)
    for i, s in enumerate(steps):
        shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
        shp.adjustments[0] = 0.35
        shp.fill.solid()
        shp.fill.fore_color.rgb = color if i < n - 1 else ORANGE
        shp.line.fill.background()
        shp.shadow.inherit = False
        shape_text(shp, [(s, 13, True, WHITE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(int(x + w + gap))


# ════════════════════════════════════════════════════════════
# S1. 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY, round_=False)
rect(s, Inches(0.9), Inches(2.55), Inches(0.14), Inches(1.6), ORANGE, round_=False)
put_text(s, Inches(1.25), Inches(2.1), Inches(11), Inches(0.5),
         [("사내 기술 공유 · PoC 4주 회고", 15, True, ORANGE, PP_ALIGN.LEFT)])
put_text(s, Inches(1.25), Inches(2.6), Inches(11.5), Inches(1.6),
         [("증권 상담원 AI 코치", 44, True, WHITE, PP_ALIGN.LEFT),
          ("“물어보면 답하고, 답하면 채점하는 AI”", 20, False, RGBColor(0xC7, 0xD2, 0xE4), PP_ALIGN.LEFT)])
put_text(s, Inches(1.25), Inches(5.6), Inches(11), Inches(0.9),
         [("무엇을 만들었나보다, 왜 그렇게 만들었나를 이야기합니다", 15, True, WHITE, PP_ALIGN.LEFT),
          ("RAG 챗봇 + 신입 훈련 시뮬레이터  |  FastAPI · React · ChromaDB · Gemini", 13, False, RGBColor(0x9C, 0xA8, 0xBC), PP_ALIGN.LEFT)])

# ════════════════════════════════════════════════════════════
# S2. 목차
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "AGENDA", "오늘 이야기 4가지")
items = [
    ("01", "무엇을 만들었나", "업무 챗봇 + 훈련 시뮬레이터, 4주 PoC"),
    ("02", "왜 그렇게 만들었나", "DB·검색 구조·스트리밍·추상화, 선택의 이유"),
    ("03", "어떤 시행착오가 있었나", "크레딧 바닥 · 사내망 SSL · PDF 표 파싱"),
    ("04", "폐쇄망에서 어떻게 쓰나", "실도입 로드맵 — 교체 포인트는 단 3곳"),
]
y = Inches(1.95)
for num, t, d in items:
    card = rect(s, Inches(0.6), y, Inches(12.13), Inches(1.08), LIGHT)
    shape_text(card, [(f"{num}   {t}", 19, True, NAVY, PP_ALIGN.LEFT),
                      (d, 13, False, GRAY, PP_ALIGN.LEFT)])
    y = Emu(int(y + Inches(1.25)))

# ════════════════════════════════════════════════════════════
# S3. 문제
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "01 · 무엇을 만들었나", "문제 — 사람이 편람을 뒤지고, 사람이 신입을 가르친다")
c1 = rect(s, Inches(0.6), Inches(2.2), Inches(5.9), Inches(3.6), LIGHT)
shape_text(c1, [("고객 응대", 15, True, GRAY, PP_ALIGN.CENTER),
                ("2~5분", 54, True, NAVY, PP_ALIGN.CENTER),
                ("숙련 상담원도 편람을 직접 검색", 14, False, INK, PP_ALIGN.CENTER),
                ("신입은 그 이상", 13, False, GRAY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
c2 = rect(s, Inches(6.85), Inches(2.2), Inches(5.9), Inches(3.6), LIGHT)
shape_text(c2, [("신입 교육(OJT)", 15, True, GRAY, PP_ALIGN.CENTER),
                ("3~6개월", 54, True, NAVY, PP_ALIGN.CENTER),
                ("선배가 1:1로 붙어야 함", 14, False, INK, PP_ALIGN.CENTER),
                ("그런데 사람이 부족함", 13, False, GRAY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
takeaway(s, "검색도, 교육도 사람의 시간을 갈아 넣는 구조 → AI로 둘 다 풀 수 있을까?")

# ════════════════════════════════════════════════════════════
# S4. 해결 — 두 가지 모드
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "01 · 무엇을 만들었나", "해결 — 하나의 데이터, 두 가지 모드")
c1 = rect(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(3.7), WHITE, line=NAVY)
shape_text(c1, [("💬 챗봇 모드", 19, True, NAVY, PP_ALIGN.LEFT),
                ("상담원이 질문하면", 13, False, GRAY, PP_ALIGN.LEFT),
                ("편람/FAQ 검색 → 즉시 답변 + 출처 표시", 15, True, INK, PP_ALIGN.LEFT),
                ("", 6, False, INK, PP_ALIGN.LEFT),
                ("“비대면 계좌 개설 절차가 뭐죠?”", 13, False, GRAY, PP_ALIGN.LEFT),
                ("→ 답변 + [출처: 편람 §계좌개설]", 13, False, ORANGE, PP_ALIGN.LEFT)])
c2 = rect(s, Inches(6.85), Inches(2.1), Inches(5.9), Inches(3.7), WHITE, line=TEAL)
shape_text(c2, [("🎓 훈련 모드", 19, True, TEAL, PP_ALIGN.LEFT),
                ("AI가 고객 역할을 맡아", 13, False, GRAY, PP_ALIGN.LEFT),
                ("질문 출제 → 신입이 답변 → AI가 채점·피드백", 15, True, INK, PP_ALIGN.LEFT),
                ("", 6, False, INK, PP_ALIGN.LEFT),
                ("AI: “양도세 언제 내요?” → 신입 답변", 13, False, GRAY, PP_ALIGN.LEFT),
                ("→ 점수 + 부족한 부분 피드백", 13, False, TEAL, PP_ALIGN.LEFT)])
takeaway(s, "같은 FAQ/편람 데이터를 — 챗봇은 “검색→답변”, 훈련은 “출제→채점”으로 다르게 쓴다")

# ════════════════════════════════════════════════════════════
# S5. 동작 방식 — RAG (왜 파인튜닝이 아닌가)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "02 · 왜 그렇게 만들었나", "왜 파인튜닝이 아니라 RAG인가")
flow(s, ["질문", "임베딩 변환", "ChromaDB\n유사 문서 검색", "LLM이 문서\n참고해 생성", "답변 + 출처"], y=Inches(2.2), h=Inches(1.15))
c = rect(s, Inches(0.6), Inches(3.85), Inches(12.13), Inches(2.4), LIGHT)
shape_text(c, [("모델은 그대로, 문서만 검색해서 보여준다 (RAG)", 16, True, NAVY, PP_ALIGN.LEFT),
               ("•  편람이 바뀌면? → 재학습 없이 재인제스트만 하면 끝", 14, False, INK, PP_ALIGN.LEFT),
               ("•  “어디서 나온 답이야?” → 검색한 문서가 곧 출처. 금융 도메인에서 필수", 14, False, INK, PP_ALIGN.LEFT),
               ("•  문서에 없으면 “모른다”고 답하게 강제 → 할루시네이션 억제", 14, False, INK, PP_ALIGN.LEFT)])
takeaway(s, "자주 바뀌는 사내 문서 + 출처 의무 → 파인튜닝보다 RAG가 구조적으로 유리")

# ════════════════════════════════════════════════════════════
# S6. 왜 ChromaDB인가
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "02 · 왜 그렇게 만들었나", "왜 벡터DB로 ChromaDB인가 — 기준은 “폐쇄망 이식성”")
opts = [
    ("Pinecone", "클라우드 전용 SaaS", "관리 불필요, 성능 좋음", "폐쇄망 반입 자체가 불가 → 탈락", False),
    ("Milvus", "대규모 분산 벡터DB", "수억 건 스케일 대응", "서버·etcd 등 설치 복잡, PoC에 과함", False),
    ("FAISS", "벡터 검색 라이브러리", "빠르고 가벼움", "DB가 아님 — 메타데이터 관리 직접 구현", False),
    ("ChromaDB", "파일 기반 임베디드 DB", "pip 한 줄 설치, 폴더 복사 = 이식 완료", "대규모엔 부적합 (PoC 규모엔 충분)", True),
]
option_cards(s, opts, y=Inches(2.15), h=Inches(4.1))
takeaway(s, "폐쇄망 반입은 “파일 복사”가 가장 강력한 배포 수단 — 서버 없는 ChromaDB 선택")

# ════════════════════════════════════════════════════════════
# S7. 왜 제목+내용 분리 (Dual Collection)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "02 · 왜 그렇게 만들었나", "검색 품질 ① — 왜 제목과 본문을 따로 임베딩하나")
put_text(s, Inches(0.9), Inches(1.6), Inches(11.8), Inches(0.5),
         [("문제: 본문이 길수록 제목의 주제 신호가 희석됨 → “양도세” 같은 짧은 질문에 취약", 14, False, GRAY, PP_ALIGN.LEFT)])
opts = [
    ("본문만 임베딩", "단일 컬렉션", "단순, 임베딩 1회", "긴 본문에서 주제 매칭 약화", False),
    ("제목+본문 합쳐서", "단일 컬렉션", "역시 단순", "가중치 조절 불가 — 본문 길수록 제목 묻힘", False),
    ("두 컬렉션 분리 + 가중 병합", "faq_titles / faq_contents", "제목:내용 비율을 파라미터로 → 그리드 서치 튜닝 가능", "임베딩 2배, ID 정합성 관리 필요", True),
]
option_cards(s, opts, y=Inches(2.25), h=Inches(3.95))
takeaway(s, "사람인HR 챗봇에서 검증된 기법(정확도 87%) — “튜닝 가능한 구조”를 사는 선택")

# ════════════════════════════════════════════════════════════
# S8. 왜 Max Pooling
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "02 · 왜 그렇게 만들었나", "검색 품질 ② — 동의어 문제와 Max Pooling")
put_text(s, Inches(0.9), Inches(1.6), Inches(11.8), Inches(0.75),
         [("“양도세” = “주식 양도 세금” → 문서 1개에 제목 여러 개를 달아 해결 (원본/유사/요약 제목)", 14, False, GRAY, PP_ALIGN.LEFT),
          ("그런데 같은 문서가 검색 결과에 여러 번 등장하면, 점수를 어떻게 합칠까?", 14, True, INK, PP_ALIGN.LEFT)])
opts = [
    ("합산 (Sum)", "등장한 제목 점수 전부 더함", "직관적", "제목 많이 단 문서가 무조건 이김 — 편향", False),
    ("평균 (Mean)", "제목 점수 평균", "합산 편향은 회피", "강하게 1번 일치한 신호가 희석됨", False),
    ("최댓값 (Max Pooling)", "가장 잘 맞은 제목 1개만 채택", "제목 개수와 무관하게 공정 — 동의어를 마음껏 추가 가능", "약한 다중 일치는 미반영 (의도된 동작)", True),
]
option_cards(s, opts, y=Inches(2.55), h=Inches(3.65))
takeaway(s, "순서도 중요: Max Pooling → 가중 병합 (순서를 바꾸면 제목 많은 문서가 부당 가중)")

# ════════════════════════════════════════════════════════════
# S9. 왜 SSE + 출처 선전송
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "02 · 왜 그렇게 만들었나", "왜 출처를 답변보다 먼저 보내나 (SSE 스트리밍)")
c1 = rect(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(3.8), LIGHT)
shape_text(c1, [("❌ 한 번에 JSON 응답", 16, True, BAD, PP_ALIGN.LEFT),
                ("LLM 생성이 다 끝나야 화면에 표시", 13, False, INK, PP_ALIGN.LEFT),
                ("→ 사용자는 5초간 빈 화면", 15, True, BAD, PP_ALIGN.LEFT),
                ("", 6, False, INK, PP_ALIGN.LEFT),
                ("상담 중인 상담원에게 5초 침묵은 치명적", 13, False, GRAY, PP_ALIGN.LEFT)])
c2 = rect(s, Inches(6.85), Inches(2.0), Inches(5.9), Inches(3.8), WHITE, line=ORANGE)
shape_text(c2, [("✅ SSE + sources 선전송", 16, True, ORANGE, PP_ALIGN.LEFT),
                ("① event: sources — 검색 끝나는 즉시 출처·신뢰도 전송", 13, False, INK, PP_ALIGN.LEFT),
                ("② event: token — 답변을 토큰 단위로 점진 표시", 13, False, INK, PP_ALIGN.LEFT),
                ("③ event: done — 명시적 종료", 13, False, INK, PP_ALIGN.LEFT),
                ("", 6, False, INK, PP_ALIGN.LEFT),
                ("출처는 검색 직후 이미 확정 — LLM을 기다릴 이유가 없다", 13, True, NAVY, PP_ALIGN.LEFT)])
takeaway(s, "“어떤 문서를 참고하는지”가 답변보다 먼저 보임 → 첫 응답 1초 + 신뢰 확보")

# ════════════════════════════════════════════════════════════
# S10. 왜 LLM 추상화
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "02 · 왜 그렇게 만들었나", "가장 잘한 결정 — LLM 서비스 추상화")
put_text(s, Inches(0.9), Inches(1.6), Inches(11.8), Inches(0.5),
         [("전제: PoC는 외부 API, 실도입은 폐쇄망 자체 호스팅 모델 → 교체는 “예정된 미래”", 14, False, GRAY, PP_ALIGN.LEFT)])
opts = [
    ("SDK 직접 호출", "rag.py 등에서 openai 직접 사용", "당장 빠름", "모델 교체 시 호출처 전부 수정", False),
    ("LangChain 도입", "프레임워크에 위임", "멀티 LLM 지원 풍부", "4주 PoC에 추상화 비용 과대, 디버깅 어려움", False),
    ("자체 인터페이스 (ABC)", "메서드 4개만: generate / generate_stream / embed / count_tokens", "교체 = 구현체 1개 추가 + DI 한 곳 수정. 테스트 모킹 쉬움", "새 모델마다 직접 구현 필요", True),
]
option_cards(s, opts, y=Inches(2.25), h=Inches(3.95))
takeaway(s, "딱 필요한 만큼만 추상화 — 이 결정이 일주일 뒤 진짜 보험금을 지급한다 (다음 장)")

# ════════════════════════════════════════════════════════════
# S11. 시행착오 ① 크레딧 바닥 → Gemini 전환
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "03 · 시행착오", "① OpenAI 크레딧이 바닥났다")
c1 = rect(s, Inches(0.6), Inches(1.95), Inches(12.13), Inches(1.5), LIGHT)
shape_text(c1, [("PoC 진행 중 API 크레딧 소진 — 충전 결제/조달 절차가 일정을 잡아먹는 상황", 15, True, INK, PP_ALIGN.LEFT),
                ("Claude는 당시 무료 임베딩 API가 제한적 → 무료 티어가 넓고 chat+embedding을 모두 주는 Gemini 선택", 13.5, False, GRAY, PP_ALIGN.LEFT)])
flow(s, ["GPT-4o\ntext-embedding-3", "GeminiService\n구현체 1개 작성", "DI 한 곳 교체", "gemini-2.5-flash\ngemini-embedding"], y=Inches(3.8), h=Inches(1.15))
c2 = rect(s, Inches(0.6), Inches(5.25), Inches(12.13), Inches(1.1), WHITE, line=ORANGE)
shape_text(c2, [("추상화의 보험금: 호출처(rag·출제·채점) 코드 변경 0줄", 15, True, ORANGE, PP_ALIGN.LEFT),
                ("단, 임베딩 차원 1536→3072 → ChromaDB는 --clear 재인제스트 필수 (공짜 점심은 없다)", 13, False, INK, PP_ALIGN.LEFT)])
takeaway(s, "교훈: 외부 의존(크레딧·결제)도 리스크다 — 추상화는 기술 부채가 아니라 보험이었다")

# ════════════════════════════════════════════════════════════
# S12. 시행착오 ② 사내망 SSL
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "03 · 시행착오", "② 사내망에서 갑자기 SSL 에러")
c1 = rect(s, Inches(0.6), Inches(1.95), Inches(12.13), Inches(1.45), LIGHT)
shape_text(c1, [("CERTIFICATE_VERIFY_FAILED: self-signed certificate in certificate chain", 14, True, BAD, PP_ALIGN.LEFT),
                ("원인: 사내 프록시(Zscaler)가 HTTPS에 자체 CA를 끼워 넣음 → SDK는 기본 CA(certifi)만 신뢰", 13.5, False, INK, PP_ALIGN.LEFT)])
c2 = rect(s, Inches(0.6), Inches(3.6), Inches(5.9), Inches(2.6), WHITE, line=BAD)
shape_text(c2, [("❌ 쉬운 길: verify=False", 15, True, BAD, PP_ALIGN.LEFT),
                ("•  검색하면 제일 먼저 나오는 해법", 13, False, INK, PP_ALIGN.LEFT),
                ("•  사내 보안 정책 위반", 13, False, INK, PP_ALIGN.LEFT),
                ("•  중간자 공격에 그대로 노출", 13, False, INK, PP_ALIGN.LEFT)])
c3 = rect(s, Inches(6.85), Inches(3.6), Inches(5.9), Inches(2.6), WHITE, line=GOOD)
shape_text(c3, [("✅ 옳은 길: 사내 CA를 신뢰 목록에", 15, True, GOOD, PP_ALIGN.LEFT),
                ("①  Windows 신뢰 저장소에서 CA 일괄 추출 (PS 스크립트화)", 13, False, INK, PP_ALIGN.LEFT),
                ("②  SSL_CERT_FILE 환경변수로 지정", 13, False, INK, PP_ALIGN.LEFT),
                ("③  REQUESTS_CA_BUNDLE 등도 동기화", 13, False, INK, PP_ALIGN.LEFT)])
takeaway(s, "폐쇄망/사내망 필수 노하우 — 스크립트로 만들어 두면 PC 교체·venv 재생성 시 재사용")

# ════════════════════════════════════════════════════════════
# S13. 시행착오 ③ PDF 표 파싱
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "03 · 시행착오", "③ 편람 PDF의 표가 검색이 안 된다")
put_text(s, Inches(0.9), Inches(1.6), Inches(11.8), Inches(0.5),
         [("업무편람 = 수수료율·코드 매핑 등 표 투성이 + 한국어 개조식 번호(1) 가. ㄱ.) — 검색 품질의 절반이 파싱에서 결정", 14, False, GRAY, PP_ALIGN.LEFT)])
opts = [
    ("pymupdf (처음 선택)", "텍스트 위주 추출", "가볍고 빠름", "표를 행/열 구조로 못 읽음", False),
    ("pdfplumber + camelot", "표 추출 조합", "표 추출 가능", "두 라이브러리 결합 복잡, 한국어 인식 부정확", False),
    ("Docling (IBM 오픈소스)", "TableFormer로 표 구조 인식", "표 구조 보존 + 헤딩 계층 → 출처 품질 향상", "패키지 무겁고 설치 오래 걸림", True),
]
option_cards(s, opts, y=Inches(2.25), h=Inches(3.3))
put_text(s, Inches(0.9), Inches(5.7), Inches(11.8), Inches(0.7),
         [("+ 한국어 개조식 깊이는 정규식으로 보강  |  표는 “전체 표” + “행 단위 문장화” 이중 적재 → 두 검색 모두 대응", 13.5, True, NAVY, PP_ALIGN.LEFT)])
takeaway(s, "교훈: RAG 정확도는 모델보다 파싱·청킹에서 먼저 갈린다")

# ════════════════════════════════════════════════════════════
# S14. 폐쇄망 실도입 로드맵
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "04 · 폐쇄망 실도입", "교체 포인트는 단 3곳 — 나머지는 그대로 들고 간다")
rows = [
    ("LLM", "Gemini 2.5 Flash (API)", "Qwen3-30B-A3B (자체 호스팅)", "오픈소스 필수. MoE 구조라 A100 1장으로 추론 가능", True),
    ("임베딩", "gemini-embedding (3072차원)", "bge-m3 (1024차원)", "오픈소스 후보. 차원이 달라져 재인제스트 필요 — 이미 한 번 해봄", True),
    ("벡터DB", "ChromaDB (파일 기반)", "ChromaDB 그대로", "처음부터 폐쇄망을 보고 고른 선택 — 폴더 복사로 이식 끝", False),
    ("백엔드/UI", "FastAPI + React", "그대로", "추상화 인터페이스 뒤에 있어 변경 없음", False),
]
y = Inches(1.95)
for name, poc, prod, why, change in rows:
    card = rect(s, Inches(0.6), y, Inches(12.13), Inches(1.02), LIGHT if not change else WHITE,
                line=TEAL if change else None)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for txt, sz, bold, col in [
        (f"{name}  ", 15, True, NAVY),
        (f"{poc}  →  {prod}", 14, True, TEAL if change else GRAY),
    ]:
        r = p.add_run(); r.text = txt
        f = r.font; f.name, f.size, f.bold = FONT, Pt(sz), bold; f.color.rgb = col
        _ea(r)
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = why
    f = r.font; f.name, f.size, f.bold = FONT, Pt(12.5), False; f.color.rgb = INK
    _ea(r)
    y = Emu(int(y + Inches(1.12)))
takeaway(s, "PoC에서 검증한 시행착오(재인제스트·CA 번들·추상화)가 그대로 폐쇄망 전환 매뉴얼이 된다", color=TEAL)

# ════════════════════════════════════════════════════════════
# S15. 우리가 남긴 것
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "마무리", "코드보다 오래 남는 것 — ‘왜’의 기록")
cards = [
    ("📐 ADR 9건", "결정마다 검토한 대안과 탈락 이유까지 기록", "“왜 ChromaDB지?”에 6개월 뒤에도 답할 수 있음"),
    ("🔧 TROUBLESHOOTING", "에러 → 원인 → 해결 → 교훈 형식으로 축적", "같은 SSL 삽질을 다음 사람이 반복하지 않음"),
    ("📝 PROMPT_HISTORY", "프롬프트 변경을 코드와 분리해 버전 기록", "정확도가 떨어지면 코드 탓인지 프롬프트 탓인지 즉시 추적"),
]
x = Inches(0.6)
for t, d1, d2 in cards:
    card = rect(s, x, Inches(2.1), Inches(3.91), Inches(3.7), LIGHT)
    shape_text(card, [(t, 17, True, NAVY, PP_ALIGN.LEFT),
                      (d1, 13, False, INK, PP_ALIGN.LEFT),
                      ("→ " + d2, 12.5, True, ORANGE, PP_ALIGN.LEFT)])
    x = Emu(int(x + Inches(4.11)))
takeaway(s, "AI가 짠 코드도 “작성자가 설명할 수 있어야 한다” — 2인 팀이 4주를 버틴 비결")

# ════════════════════════════════════════════════════════════
# S16. 클로징
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY, round_=False)
put_text(s, Inches(1.25), Inches(2.0), Inches(11), Inches(3.5),
         [("오늘의 세 줄 요약", 18, True, ORANGE, PP_ALIGN.LEFT),
          ("1.  기술 선택의 기준은 처음부터 “폐쇄망 이식성”이었다 — ChromaDB·추상화·오픈소스 후보", 17, True, WHITE, PP_ALIGN.LEFT),
          ("2.  시행착오는 기록하면 자산이 된다 — 크레딧·SSL·PDF 파싱 전부 문서로", 17, True, WHITE, PP_ALIGN.LEFT),
          ("3.  RAG 품질은 모델이 아니라 구조에서 나온다 — 분리 임베딩·Max Pooling·파싱", 17, True, WHITE, PP_ALIGN.LEFT)])
put_text(s, Inches(1.25), Inches(5.9), Inches(11), Inches(0.8),
         [("Q&A  |  자세한 근거는 리포지토리 docs/adr/ 에 전부 있습니다", 15, False, RGBColor(0xC7, 0xD2, 0xE4), PP_ALIGN.LEFT)])

# ── 저장 ────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AI코치_PoC_공유발표.pptx")
prs.save(out)
print(f"저장 완료: {out}  (슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장)" if False else f"저장 완료: {out}")
