# -*- coding: utf-8 -*-
"""HICT_PPT.pptx 고객사 설명회용 재편집 엔진.

기존 손제작 덱(18장)의 디자인 토큰을 그대로 복제해, 신규 6장(시연/기술난점x4/시사점/CTA/부록)을
무봉제로 삽입하고 1·2·3장 톤을 '고객사 동행 제안'으로 전환한다.

설계 결정(api-spec.md 섹션 없음 — 발표자료라 별도):
- 기존 HICT_PPT.pptx는 python-pptx 절대좌표 손제작 덱. 생성 스크립트 부재 → 직접 XML 편집.
- 카피 텍스트는 워크플로(hict-ppt-copy)가 만든 slides_copy.json 에서 읽는다.
- 슬라이드 순서 재배치는 sldIdLst 조작(_move_slide)으로 처리.

연관: 측정한 디자인 토큰은 /tmp/design_tokens.txt 참조. 빌더는 make_ppt_hanwha_poc.py 헬퍼 패턴 차용.
실행: python edit_hict_ppt.py  (산출물: HICT_PPT.pptx 덮어쓰기, 사전 백업됨)
"""
import copy
import io
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── 디자인 토큰 (HICT_PPT.pptx 실측값) ──────────────────────────
ORANGE = RGBColor(0xFF, 0x66, 0x00)   # kicker 번호, 강조
INK = RGBColor(0x2A, 0x2A, 0x2A)      # 본문 제목, 다크 바
GRAY = RGBColor(0x56, 0x56, 0x56)     # 카드 본문
FOOT = RGBColor(0x8C, 0x8C, 0x89)     # 푸터
WMARK = RGBColor(0x56, 0x56, 0x56)    # 상단 워터마크
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)  # 카드 흰 배경
CARD_HL = RGBColor(0xFF, 0xF6, 0xEF)  # 카드 연오렌지 강조
CARD_LINE = RGBColor(0xD6, 0xD6, 0xD2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOOD = RGBColor(0x15, 0x80, 0x3D)
FONT = "맑은 고딕"

# 고객사(타 증권사) 설명회용 — '선제안' 톤 제거, 사례 공유 톤
WMARK_TEXT = "한화시스템 · 2026 상반기 AX 추진 사례 공유"
WMARK_OLD = "한화투자증권 RAG AI 상담 어시스턴트 AX 선제안"  # 기존 슬라이드 일괄 치환 대상
COPYRIGHT = "ⓒ Hanwha Systems. All rights reserved"


def _set_font(run, size=None, color=None, bold=None, name=FONT):
    f = run.font
    f.name = name
    # 한글 폰트 강제(eastasian)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color


def text_box(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, line_spacing=None):
    """runs = [(text, size, color, bold), ...] 여러 run을 한 문단에. 줄바꿈은 '\n' 분리 run로."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for (txt, size, color, bold) in runs:
        if txt == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            if line_spacing:
                p.line_spacing = line_spacing
            continue
        r = p.add_run()
        r.text = txt
        _set_font(r, size=size, color=color, bold=bold)
    return tb


def card(slide, x, y, w, h, fill=CARD_BG, line=CARD_LINE, line_w=1.0, round_=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    # 라운드 정도 살짝 작게
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


def add_header(slide, kicker_no, kicker_label, title, page_no):
    """좌상단 kicker 번호+라벨, 타이틀, 상단 워터마크, 우하단 푸터·페이지."""
    # 상단 워터마크 (마스터 로고 left=10.58 이전에 끝나도록 w 제한)
    text_box(slide, 1.54, 0.48, 6.5, 0.18, [(WMARK_TEXT, 8.5, WMARK, False)])
    # kicker 번호 (큰 오렌지)
    text_box(slide, 1.08, 0.66, 0.55, 0.50, [(kicker_no, 30, ORANGE, True)])
    # 타이틀 (w=9.0 — 마스터 로고 left=10.58과 겹치지 않도록)
    runs = []
    if kicker_label:
        runs.append((kicker_label + ": ", 20, ORANGE, True))
    runs.append((title, 20, INK, True))
    text_box(slide, 1.54, 0.74, 9.0, 0.55, runs, anchor=MSO_ANCHOR.MIDDLE)
    # 푸터
    text_box(slide, 0.75, 7.12, 3.2, 0.16, [(COPYRIGHT, 7.5, FOOT, False)])
    text_box(slide, 6.24, 7.12, 0.85, 0.16, [("- %d -" % page_no, 7.5, FOOT, False)],
             align=PP_ALIGN.CENTER)


def takeaway_bar(slide, text, y=5.95):
    bar = card(slide, 1.10, y, 11.10, 0.55, fill=INK, line=WHITE, round_=True)
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    _set_font(r, size=11.5, color=WHITE, bold=True)


def note_line(slide, text, y=6.70):
    text_box(slide, 1.10, y, 11.10, 0.30,
             [("쉬운 풀이 — " + text if not text.startswith("쉬운") else text, 9.5, GRAY, False)])


# ── 카드 그리드: N개 카드를 가로로 균등 배치 ──────────────────
def card_row(slide, items, y=2.00, h=3.55, top=1.95, highlight_idx=1):
    """items = [{num, head, desc, metric}]. 가로 균등. highlight_idx 카드만 연오렌지."""
    n = len(items)
    margin = 0.88
    gap = 0.30
    total_w = 13.333 - margin * 2
    cw = (total_w - gap * (n - 1)) / n
    # 카드가 좁을수록(많을수록) 폰트·간격 축소해 desc 공간 확보
    compact = n >= 5
    pad = 0.18 if compact else 0.25
    head_sz = 12 if compact else 14.5
    metric_sz = 11 if compact else 15
    desc_sz = 8.6 if compact else 9.8
    head_h = 0.50 if compact else 0.60
    metric_h = 0.26 if compact else 0.34
    metric_gap = 0.34 if compact else 0.46
    for i, it in enumerate(items):
        x = margin + i * (cw + gap)
        fill = CARD_HL if i == highlight_idx else CARD_BG
        card(slide, x, y, cw, h, fill=fill)
        cy = y + (0.22 if compact else 0.30)
        if it.get('num'):
            text_box(slide, x + pad, cy, 0.6, 0.24, [(it['num'], 12 if compact else 13, ORANGE, True)])
            cy += 0.34 if compact else 0.42
        text_box(slide, x + pad, cy, cw - pad * 2, head_h,
                 [(it['head'], head_sz, INK, True)], line_spacing=1.0)
        cy += head_h + (0.02 if compact else 0.02)
        if it.get('metric'):
            text_box(slide, x + pad, cy, cw - pad * 2, metric_h, [(it['metric'], metric_sz, ORANGE, True)])
            cy += metric_gap
        text_box(slide, x + pad, cy, cw - pad * 2, h - (cy - y) - 0.18,
                 [(it['desc'], desc_sz, GRAY, False)], line_spacing=1.04)


# ── 카드 그리드: 2열 N행 (난점→해결 처럼 항목 많을 때) ─────────
def card_grid_2col(slide, items, y=1.95, row_h=1.05, gap=0.22):
    margin = 0.82
    gap_x = 0.30
    cw = (13.333 - margin * 2 - gap_x) / 2
    for i, it in enumerate(items):
        col = i % 2
        row = i // 2
        x = margin + col * (cw + gap_x)
        cy = y + row * (row_h + gap)
        fill = CARD_HL if it.get('hl') else CARD_BG
        card(slide, x, cy, cw, row_h, fill=fill)
        text_box(slide, x + 0.22, cy + 0.15, 0.48, 0.22,
                 [(it.get('num', ''), 11, ORANGE, True)])
        text_box(slide, x + 0.66, cy + 0.13, cw - 0.85, 0.28,
                 [(it['head'], 12.5, INK, True)], line_spacing=1.0)
        text_box(slide, x + 0.22, cy + 0.46, cw - 0.42, row_h - 0.54,
                 [(it['desc'], 9.2, GRAY, False)], line_spacing=1.02)


# ── 빈 블랭크 슬라이드 추가 ────────────────────────────────────
def add_blank(prs):
    # 레이아웃 6 = 완전 빈 화면 (기존 덱과 동일 가정). 없으면 마지막 레이아웃.
    layouts = prs.slide_layouts
    blank = None
    for lo in layouts:
        if len(lo.placeholders) == 0:
            blank = lo
            break
    if blank is None:
        blank = layouts[-1]
    return prs.slides.add_slide(blank)


# ── 슬라이드 순서 이동 ─────────────────────────────────────────
def move_slide(prs, from_idx, to_idx):
    """0-based. from_idx 슬라이드를 to_idx 위치로 이동."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    el = ids[from_idx]
    sldIdLst.remove(el)
    ids2 = list(sldIdLst)
    if to_idx >= len(ids2):
        sldIdLst.append(el)
    else:
        sldIdLst.insert(to_idx, el)


if __name__ == '__main__':
    print("이 파일은 빌더 모듈입니다. build_slides.py 에서 import 해 사용합니다.")
