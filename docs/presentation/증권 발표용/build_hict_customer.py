# -*- coding: utf-8 -*-
"""HICT_PPT.pptx → 고객사 설명회용 24장 덱 빌드.

입력:  HICT_PPT.pptx(18장, 손제작) + /tmp/slides_final.json(워크플로 카피)
출력:  HICT_PPT.pptx 덮어쓰기 (사전 백업 완료)

전략:
- 1·2·3장(타이틀/목차/Exec)은 기존 도형을 전부 지우고 톤전환 카피로 재구성.
- 신규 6장(시연/기술난점4/시사점/CTA) + 부록 2장은 빈 슬라이드로 추가 후 move_slide로 제자리 삽입.
- 기존 4~18장은 그대로 두되, 일부 캡션을 '귀사 적용' 관점으로 손질(별도 함수).

연관: edit_hict_ppt.py(빌더 헬퍼), 디자인 토큰 /tmp/design_tokens.txt.
api-spec 해당 없음(발표자료). 슬라이드 순서는 sldIdLst 조작.
"""
import io
import json
import os
import sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import edit_hict_ppt as B  # 빌더 헬퍼
from edit_hict_ppt import (ORANGE, INK, GRAY, WHITE, CARD_BG, CARD_HL,
                           add_blank, add_header, card, card_row, card_grid_2col,
                           move_slide, note_line, takeaway_bar, text_box)

HERE = os.path.dirname(os.path.abspath(__file__))
# 입력은 깨끗한 백업, 출력은 별도 파일(PowerPoint가 원본을 잠그고 있을 수 있음).
SRC = os.environ.get("HICT_SRC",
                     os.path.join(HERE, "HICT_PPT_backup_고객사재편집전_20260630_153844.pptx"))
PPTX = os.environ.get("HICT_OUT", os.path.join(HERE, "HICT_PPT_고객사판.pptx"))
COPY = json.load(io.open("/tmp/slides_final.json", encoding="utf-8"))
C = {s["id"]: s for s in COPY}


def clear_slide(slide):
    """슬라이드의 모든 도형 제거(배경/레이아웃 placeholder는 유지)."""
    sp_tree = slide.shapes._spTree
    for shp in list(slide.shapes):
        sp_tree.remove(shp._element)


# ────────────────────────────────────────────────────────────
# 톤전환: 1·2·3장 재구성
# ────────────────────────────────────────────────────────────
def build_title(slide):
    s = C["title"]
    clear_slide(slide)
    # 상단 작은 kicker 라벨
    text_box(slide, 1.10, 1.55, 11.0, 0.30, [(s["kicker"], 13, ORANGE, True)])
    # 메인 타이틀(큰)
    text_box(slide, 1.10, 2.15, 11.2, 1.30, [(s["title"], 34, INK, True)],
             line_spacing=1.05)
    # 부제
    text_box(slide, 1.10, 3.55, 11.0, 0.70, [(s["subtitle"], 15, GRAY, False)],
             line_spacing=1.15)
    # 칩 3개(body)
    chips = s["body"]
    n = len(chips)
    cw, gap, x0, y = 3.55, 0.30, 1.10, 4.70
    for i, c in enumerate(chips):
        x = x0 + i * (cw + gap)
        fill = CARD_HL if i == 1 else CARD_BG
        card(slide, x, y, cw, 1.35, fill=fill)
        text_box(slide, x + 0.25, y + 0.22, cw - 0.5, 0.28,
                 [(c["head"], 13.5, ORANGE, True)])
        text_box(slide, x + 0.25, y + 0.62, cw - 0.5, 0.65,
                 [(c["desc"], 9.3, GRAY, False)], line_spacing=1.05)
    # 날짜/카피
    text_box(slide, 1.10, 6.55, 6.0, 0.22, [("2026.06", 11, GRAY, True)])
    text_box(slide, 0.75, 7.12, 4.5, 0.16, [(B.COPYRIGHT, 7.5, B.FOOT, False)])
    if s.get("note"):
        text_box(slide, 1.10, 6.85, 11.0, 0.30, [(s["note"], 8.2, B.FOOT, False)])


def build_toc(slide, page_no):
    s = C["toc"]
    clear_slide(slide)
    add_header(slide, "", s["kicker"], s["title"], page_no)
    text_box(slide, 1.10, 1.45, 9.5, 0.28, [(s["subtitle"], 11.5, GRAY, False)])
    # 8항목 2열 4행 (gap=0.15로 타이트하게, row_h=1.10으로 약간 더 여유)
    items = [{"num": b["num"], "head": b["head"], "desc": b["desc"]} for b in s["body"]]
    card_grid_2col(slide, items, y=1.88, row_h=1.10, gap=0.15)
    # 카드 4행 끝 = 1.88 + 4*(1.10+0.15) - 0.15 = 6.73 → bar는 6.78
    takeaway_bar(slide, s["takeaway"], y=6.82)


def build_exec(slide, page_no):
    s = C["exec"]
    clear_slide(slide)
    add_header(slide, "", s["kicker"], s["title"], page_no)
    text_box(slide, 1.10, 1.50, 9.5, 0.28, [(s["subtitle"], 11.5, GRAY, False)])
    items = [{"head": b["head"], "desc": b["desc"], "metric": b.get("metric")}
             for b in s["body"]]
    # 카드 y=1.90, h=3.60 → 끝 = 5.50. takeaway = 5.65
    card_row(slide, items, y=1.90, h=3.60, highlight_idx=1)
    takeaway_bar(slide, s["takeaway"], y=5.65)
    # note_line 제거 — 카드에 충분한 설명 포함, 쉬운풀이 바가 겹치는 문제 해결


# ────────────────────────────────────────────────────────────
# 신규 슬라이드 빌더(공통: 카드행 또는 2열그리드)
# ────────────────────────────────────────────────────────────
def build_cardrow_slide(slide, sid, page_no, kicker_label_split=True, highlight_idx=1):
    s = C[sid]
    add_header(slide, "", s["kicker"], s["title"], page_no)
    y0 = 1.50
    if s.get("subtitle"):
        text_box(slide, 1.10, y0, 11.2, 0.40, [(s["subtitle"], 11, GRAY, False)],
                 line_spacing=1.1)
        y0 += 0.45
    items = [{"num": b.get("num"), "head": b["head"], "desc": b["desc"],
              "metric": b.get("metric")} for b in s["body"]]
    n = len(items)
    # 카드 높이/타이틀 길이에 따라 동적 조정
    card_h = 3.55 if n >= 5 else (3.40 if n == 4 else 3.55)
    card_row(slide, items, y=y0 + 0.05, h=card_h, highlight_idx=highlight_idx)
    ty = 5.72 if n <= 3 else 5.68
    takeaway_bar(slide, s["takeaway"], y=ty)
    if s.get("note"):
        note_line(slide, s["note"], y=ty + 0.62)


def build_demo(slide, page_no, video_path=None, poster_path=None):
    """시연 슬라이드. 좌측 = 영상(한 개, 응대→AI코치 연속), 우측 = 2모드 설명 카드.
    video_path 가 주어지면 PowerPoint 비디오로 임베드(MP4/H.264 권장)."""
    s = C["demo"]
    add_header(slide, "", s["kicker"], s["title"], page_no)
    text_box(slide, 1.10, 1.48, 11.2, 0.40, [(s["subtitle"], 11, GRAY, False)],
             line_spacing=1.1)

    # ── 좌측: 영상 영역 (1280x800 비율 = 16:10) ──
    vx, vy, vw, vh = 1.10, 2.05, 7.10, 4.44  # 16:10 → 7.10 x 4.44
    if video_path and os.path.exists(video_path):
        try:
            # 프레임 배경(영상 테두리)
            card(slide, vx - 0.06, vy - 0.06, vw + 0.12, vh + 0.12,
                 fill=B.RGBColor(0x1A, 0x1A, 0x1A), line=None)
            mv = slide.shapes.add_movie(
                video_path, Inches(vx), Inches(vy), Inches(vw), Inches(vh),
                poster_frame_image=poster_path, mime_type="video/mp4")
            # 슬라이드 진입 시 자동재생하도록 설정(없으면 클릭 재생)
            _set_video_autoplay(mv)
            placed = True
        except Exception as e:
            print("  [경고] 영상 임베드 실패, 플레이스홀더로 대체:", e)
            placed = False
    else:
        placed = False
    if not placed:
        card(slide, vx, vy, vw, vh, fill=B.RGBColor(0xF3, 0xF4, 0xF6),
             line=B.RGBColor(0xBF, 0xBF, 0xBF))
        text_box(slide, vx + 0.3, vy + vh / 2 - 0.2, vw - 0.6, 0.4,
                 [("▶  시연 영상 자리 — 영상을 이 박스 위에 끌어다 놓으세요", 11, GRAY, True)],
                 align=PP_ALIGN.CENTER)

    # ── 우측: 2모드 설명 카드 ──
    cards = s["body"]
    rx, rw = 8.45, 3.75
    ry, rh, rgap = 2.05, 2.10, 0.24
    for i, c in enumerate(cards):
        cy = ry + i * (rh + rgap)
        card(slide, rx, cy, rw, rh, fill=CARD_HL if i == 0 else CARD_BG)
        text_box(slide, rx + 0.24, cy + 0.18, 0.6, 0.24, [(c["num"], 12.5, ORANGE, True)])
        text_box(slide, rx + 0.74, cy + 0.16, rw - 1.0, 0.50,
                 [(c["head"], 12.5, INK, True)], line_spacing=1.02)
        text_box(slide, rx + 0.24, cy + 0.72, rw - 0.5, rh - 1.05,
                 [(c["desc"], 9.0, GRAY, False)], line_spacing=1.06)
        if c.get("metric"):
            text_box(slide, rx + 0.24, cy + rh - 0.34, rw - 0.5, 0.26,
                     [(c["metric"], 9.0, ORANGE, True)])

    takeaway_bar(slide, s["takeaway"], y=5.95)
    note_txt = ("재생: 슬라이드쇼에서 자동재생(또는 영상 클릭). 좌측 한 영상에 응대→AI코치 모드가 이어집니다."
                if placed else
                "📎 " + s.get("note", ""))
    text_box(slide, 1.10, 6.62, 11.10, 0.45, [(note_txt, 8.6, B.FOOT, False)],
             line_spacing=1.05)


def _set_video_autoplay(movie_shape):
    """영상이 슬라이드 진입 시 자동재생되도록 XML 타이밍 설정.
    python-pptx는 기본 '클릭 재생'으로 넣으므로 timing 노드를 자동재생으로 교체."""
    from pptx.oxml.ns import qn
    # movie_shape._element 는 <p:pic>. 같은 슬라이드 timing(<p:timing>)에 자동재생 cond 삽입.
    pic = movie_shape._element
    # videoFile rId / shape id
    nv = pic.find(qn('p:nvPicPr'))
    if nv is None:
        return
    cNvPr = nv.find(qn('p:cNvPr'))
    if cNvPr is None:
        return
    sp_id = cNvPr.get('id')
    slide_el = pic.getparent().getparent()  # spTree -> cSld -> ... 실제론 sld
    # sld 루트 탐색
    root = pic
    while root.tag != qn('p:sld') and root.getparent() is not None:
        root = root.getparent()
    if root.tag != qn('p:sld'):
        return
    # 기존 timing 제거 후 자동재생 timing 삽입.
    # 구조: tnLst > par(tmRoot) > seq(mainSeq) > par > par > par(clickEffect로 보이나
    #       cond delay=0 → 슬라이드 진입 즉시) > cmd playFrom(0). 모든 여는 태그를 정확히 닫음.
    for old in root.findall(qn('p:timing')):
        root.remove(old)
    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    timing_xml = f'''<p:timing xmlns:p="{P}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetClass="mediacall" presetID="1" fill="hold" nodeType="afterEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:cmd type="call" cmd="playFrom(0.0)">
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="1" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
                                    </p:cBhvr>
                                  </p:cmd>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''
    from pptx.oxml import parse_xml
    root.append(parse_xml(timing_xml))


def build_grid_slide(slide, sid, page_no, row_h=1.35):
    """tech-wall, appendix처럼 2~4 큰 카드를 2열로."""
    s = C[sid]
    add_header(slide, "", s["kicker"], s["title"], page_no)
    y0 = 1.48
    if s.get("subtitle"):
        text_box(slide, 1.10, y0, 11.2, 0.42, [(s["subtitle"], 11, GRAY, False)],
                 line_spacing=1.1)
        y0 += 0.50
    items = s["body"]
    n = len(items)
    margin, gap_x = 0.82, 0.30
    cw = (13.333 - margin * 2 - gap_x) / 2
    # 카드 높이: 2개면 크게, 4개면 2열2행
    if n <= 2:
        ch = 3.30
        for i, c in enumerate(items):
            x = margin + i * (cw + gap_x)
            card(slide, x, y0, cw, ch, fill=CARD_HL if i == 0 else CARD_BG)
            cyy = y0 + 0.25
            if c.get("num"):
                text_box(slide, x + 0.28, cyy, 0.7, 0.26, [(c["num"], 13, ORANGE, True)])
                cyy += 0.42
            text_box(slide, x + 0.28, cyy, cw - 0.55, 0.55, [(c["head"], 13, INK, True)],
                     line_spacing=1.05)
            cyy += 0.62
            text_box(slide, x + 0.28, cyy, cw - 0.55, ch - (cyy - y0) - 0.55,
                     [(c["desc"], 9.6, GRAY, False)], line_spacing=1.12)
            if c.get("metric"):
                text_box(slide, x + 0.28, y0 + ch - 0.42, cw - 0.55, 0.30,
                         [(c["metric"], 9.6, ORANGE, True)])
    else:
        rows = (n + 1) // 2
        ch = 3.30 / rows if rows > 1 else 3.0
        ch = max(ch, 1.50)
        for i, c in enumerate(items):
            col, row = i % 2, i // 2
            x = margin + col * (cw + gap_x)
            cy = y0 + row * (ch + 0.18)
            card(slide, x, cy, cw, ch, fill=CARD_HL if c.get("hl") else CARD_BG)
            text_box(slide, x + 0.24, cy + 0.14, 0.6, 0.22,
                     [(c.get("num", ""), 11.5, ORANGE, True)])
            text_box(slide, x + 0.74, cy + 0.12, cw - 1.0, 0.26,
                     [(c["head"], 11.8, INK, True)])
            text_box(slide, x + 0.24, cy + 0.46, cw - 0.50, ch - 0.85,
                     [(c["desc"], 8.8, GRAY, False)], line_spacing=1.04)
            if c.get("metric"):
                text_box(slide, x + 0.24, cy + ch - 0.30, cw - 0.50, 0.24,
                         [(c["metric"], 8.8, ORANGE, True)])
    ty = 6.05
    takeaway_bar(slide, s["takeaway"], y=ty)
    if s.get("note"):
        note_line(slide, s["note"], y=ty + 0.62)


# ────────────────────────────────────────────────────────────
# 메인
# ────────────────────────────────────────────────────────────
def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    assert n0 == 18, "기존 덱이 18장이 아님: %d" % n0

    # 1) 톤전환: 기존 1·2·3장 in-place 재구성
    build_title(prs.slides[0])
    build_toc(prs.slides[1], 2)
    build_exec(prs.slides[2], 3)

    # 2) 신규 슬라이드를 끝에 추가한 뒤 move로 제자리 삽입.
    #    최종 순서(0-based 목표 인덱스):
    #    0 타이틀 / 1 목차 / 2 exec / 3 제안배경(기존4) /
    #    4 솔루션개요(기존5) / 5 상담흐름(기존6) / 6 [신규 demo] /
    #    7 기술구조(기존7) / 8 검색정확도(기존8) /
    #    9 [tech-overview] 10[wall1] 11[wall2] 12[wall3] /
    #    13 [insights] /
    #    14 PoC근거(기존9) / 15 실도입아키(기존10) / 16 보안운영(기존11) /
    #    17 비전(기존12) / 18 기대효과ROI(기존13) / 19 ROI산정(기존14) /
    #    20 로드맵(기존15) / 21 운영체계(기존16) / 22 성공기준(기존17) /
    #    23 [cta] / 24 Q&A(기존18) / 25 [appendix-search] 26 [appendix-wiki]
    #
    # 추가 순서대로 빌드(맨 끝에 붙음), 그 뒤 한꺼번에 재정렬.
    # 시연 영상: 오늘 생성된 MP4(H.264) — PowerPoint 임베드 호환. webm(VP8)은 재생 불가.
    video = os.environ.get("HICT_VIDEO", os.path.join(
        HERE, "..", "..", "..", "demo_시연영상_실제버전.mp4"))
    video = os.path.normpath(video)
    poster = os.path.join(HERE, "demo_poster.png")
    if not os.path.exists(poster):
        poster = None

    new_specs = [
        ("demo", lambda sl, pg: build_demo(sl, pg, video_path=video, poster_path=poster)),
        ("tech-overview", lambda sl, pg: build_cardrow_slide(sl, "tech-overview", pg, highlight_idx=4)),
        ("tech-wall-1", lambda sl, pg: build_grid_slide(sl, "tech-wall-1", pg)),
        ("tech-wall-2", lambda sl, pg: build_grid_slide(sl, "tech-wall-2", pg)),
        ("tech-wall-3", lambda sl, pg: build_grid_slide(sl, "tech-wall-3", pg)),
        ("insights", lambda sl, pg: build_cardrow_slide(sl, "insights", pg, highlight_idx=0)),
        ("cta", lambda sl, pg: build_grid_slide(sl, "cta", pg)),
        ("appendix-search", lambda sl, pg: build_grid_slide(sl, "appendix-search", pg)),
        ("appendix-wiki", lambda sl, pg: build_grid_slide(sl, "appendix-wiki", pg)),
    ]
    new_slides = {}
    for sid, fn in new_specs:
        sl = add_blank(prs)
        new_slides[sid] = sl
        fn(sl, 0)  # 페이지번호는 마지막에 일괄 재부여

    # 현재 sldIdLst 순서: [0..17 기존][18 demo,19 t-ov,20 w1,21 w2,22 w3,23 ins,24 cta,25 apA,26 apB]
    # 목표 인덱스로 이동. move_slide는 한 번에 하나씩 — 뒤에서부터 안전하게 배치.
    # 각 신규 슬라이드의 '현재 element'를 직접 추적해 위치 지정.
    def idx_of(slide):
        return list(prs.slides).index(slide)

    # 삽입 목표(작은 인덱스부터 처리하되, 이동 후 재계산)
    targets = [
        (new_slides["demo"], 6),
        (new_slides["tech-overview"], 9),
        (new_slides["tech-wall-1"], 10),
        (new_slides["tech-wall-2"], 11),
        (new_slides["tech-wall-3"], 12),
        (new_slides["insights"], 13),
        (new_slides["cta"], 23),
        (new_slides["appendix-search"], 25),
        (new_slides["appendix-wiki"], 26),
    ]
    for slide, tgt in targets:
        cur = idx_of(slide)
        move_slide(prs, cur, tgt)

    # 3) 기존 슬라이드 상단 워터마크 '...AX 선제안' → '사례 공유' 톤으로 일괄 치환
    _fix_watermarks(prs)
    # 3-1) 전체 '귀사' → '고객사' 일괄 치환
    _fix_guisa(prs)

    # 4) 페이지번호 일괄 재부여(신규 슬라이드 + 기존 푸터 페이지 텍스트 갱신)
    _renumber(prs)

    prs.save(PPTX)
    print("저장 완료:", PPTX, "| 총", len(prs.slides), "장")


def _fix_guisa(prs):
    """모든 슬라이드에서 '귀사' 를 '고객사' 로 치환 (격식 있되 한 단계 낮춤).
    단, '귀사와 함께 할 일' 등 어색해지는 표현은 더 자연스럽게 재작성."""
    PAIRS = [
        ("귀사와 함께 할 일", "함께 할 일"),
        ("귀사도 동일하고", "고객사도 동일하고"),
        ("귀사가 마주칠", "고객사가 마주칠"),
        ("귀사의", "고객사의"),
        ("귀사에", "고객사에"),
        ("귀사 AX", "고객사 AX"),
        ("귀사 데이터", "고객사 데이터"),
        ("귀사 환경", "고객사 환경"),
        ("귀사 인력", "고객사 인력"),
        ("귀사 과제", "고객사 과제"),
        ("귀사", "고객사"),  # 나머지 잔존 처리
    ]
    for slide in prs.slides:
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in PAIRS:
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def _fix_watermarks(prs):
    """모든 슬라이드의 기존 워터마크 문구를 고객사 설명회용으로 치환."""
    for slide in prs.slides:
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            if B.WMARK_OLD in shp.text_frame.text:
                for para in shp.text_frame.paragraphs:
                    if para.runs and B.WMARK_OLD in "".join(r.text for r in para.runs):
                        para.runs[0].text = B.WMARK_TEXT
                        for r in para.runs[1:]:
                            r.text = ""


def _renumber(prs):
    """모든 슬라이드 우하단 '- N -' 페이지번호를 실제 위치로 갱신.
    기존 슬라이드는 텍스트박스 안 '- N -' 패턴을 찾아 치환, 신규는 add_header가 이미 0으로 박았으니 재설정."""
    import re
    for i, slide in enumerate(prs.slides, start=1):
        found = False
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            t = shp.text_frame.text.strip()
            if re.fullmatch(r"-\s*\d+\s*-", t) or t == "- 0 -":
                for para in shp.text_frame.paragraphs:
                    if para.runs:
                        para.runs[0].text = "- %d -" % i
                        for r in para.runs[1:]:
                            r.text = ""
                        break
                found = True
                break
        if not found and i > 1:
            # 페이지번호 텍스트박스가 없는 슬라이드(예: Q&A)엔 새로 추가
            text_box(slide, 6.24, 7.12, 0.85, 0.16,
                     [("- %d -" % i, 7.5, B.FOOT, False)], align=PP_ALIGN.CENTER)


if __name__ == "__main__":
    main()
