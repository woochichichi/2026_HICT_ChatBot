# -*- coding: utf-8 -*-
"""한화투자증권 PoC 공유회용 30분 발표자료(.pptx) 생성기 — v2.

개선점(2026-06-20):
- 비IT 청중 배려: 어려운 용어는 슬라이드 하단 '쉬운 풀이' 각주 + 마지막 '용어 사전' 1장.
- '한눈에 보는 전체 그림' 슬라이드를 앞에 배치해 큰 그림을 먼저 이해.
- 최신 기능 반영: 응대/AI코치/제보 3기능, 시나리오 뱅크·커리큘럼·복습, 페르소나 8종, TTS,
  출처 위키링크. 쉬운 비유(RAG=오픈북 시험 등)로 서술.

실행:  .venv/Scripts/python.exe docs/presentation/make_ppt_hanwha_poc.py
산출물: docs/presentation/한화투자증권_AI상담_PoC_30min.pptx
로고 없음(요청). 헬퍼는 make_ppt.py 패턴 복사 + 한화 오렌지.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 한화 디자인 토큰 ─────────────────────────────────────────
ORANGE = RGBColor(0xF3, 0x73, 0x21)
ORANGED = RGBColor(0xD8, 0x5C, 0x12)
NAVY = RGBColor(0x13, 0x29, 0x4B)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
INK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
SOFT = RGBColor(0xFF, 0xF3, 0xEC)
GOOD = RGBColor(0x15, 0x80, 0x3D)
BAD = RGBColor(0xB9, 0x1C, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
_PAGE = [0]


# ── 헬퍼 ─────────────────────────────────────────────────────
def _ea(run):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def put_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name, f.size, f.bold = FONT, Pt(size), bold
        f.color.rgb = color
        _ea(run)
    return box


def rect(slide, x, y, w, h, fill, line=None, round_=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if round_:
        try:
            shape.adjustments[0] = 0.06
        except Exception:
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def shape_text(shape, lines, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name, f.size, f.bold = FONT, Pt(size), bold
        f.color.rgb = color
        _ea(run)


def slide():
    _PAGE[0] += 1
    return prs.slides.add_slide(BLANK)


def header(s, kicker, title):
    rect(s, Inches(0.6), Inches(0.42), Inches(0.12), Inches(0.95), ORANGE, round_=False)
    put_text(s, Inches(0.9), Inches(0.35), Inches(11.8), Inches(0.4),
             [(kicker, 13, True, ORANGE, PP_ALIGN.LEFT)])
    put_text(s, Inches(0.9), Inches(0.68), Inches(11.8), Inches(0.78),
             [(title, 26, True, NAVY, PP_ALIGN.LEFT)])
    put_text(s, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
             [(str(_PAGE[0]), 10, False, GRAY, PP_ALIGN.RIGHT)])


def takeaway(s, text, color=NAVY, y=Inches(6.5)):
    bar = rect(s, Inches(0.6), y, Inches(12.13), Inches(0.62), color)
    shape_text(bar, [(text, 14.5, True, WHITE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)


def note(s, text, y=Inches(6.95)):
    """일반 각주(회색, 작게)."""
    put_text(s, Inches(0.9), y, Inches(11.4), Inches(0.4),
             [(text, 9.5, False, GRAY, PP_ALIGN.LEFT)])


def terms(s, pairs, y=Inches(6.25)):
    """비IT용 '쉬운 풀이' 각주 박스. pairs: [(용어, 쉬운설명), ...]"""
    box = rect(s, Inches(0.6), y, Inches(12.13), Inches(0.95), SOFT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    p.space_after = Pt(2)
    r = p.add_run(); r.text = "쉬운 풀이  "
    f = r.font; f.name, f.size, f.bold = FONT, Pt(10.5), True; f.color.rgb = ORANGED; _ea(r)
    for term, desc in pairs:
        rp = tf.add_paragraph(); rp.space_after = Pt(1)
        r1 = rp.add_run(); r1.text = f"· {term}: "
        f = r1.font; f.name, f.size, f.bold = FONT, Pt(11), True; f.color.rgb = NAVY; _ea(r1)
        r2 = rp.add_run(); r2.text = desc
        f = r2.font; f.name, f.size, f.bold = FONT, Pt(11), False; f.color.rgb = INK; _ea(r2)


def option_cards(s, options, y=Inches(2.3), h=Inches(3.9)):
    n = len(options)
    gap = Inches(0.3)
    total = Inches(12.13)
    w = Emu(int((total - gap * (n - 1)) / n))
    x = Inches(0.6)
    for name, desc, pro, con, chosen in options:
        card = rect(s, x, y, w, h, WHITE if chosen else LIGHT, line=ORANGE if chosen else None)
        shape_text(card, [
            (("✅ " if chosen else "") + name, 15, True, ORANGE if chosen else INK, PP_ALIGN.LEFT),
            (desc, 11.5, False, GRAY, PP_ALIGN.LEFT),
            ("👍 " + pro, 12, False, GOOD, PP_ALIGN.LEFT),
            ("👎 " + con, 12, False, BAD, PP_ALIGN.LEFT),
        ])
        x = Emu(int(x + w + gap))


def flow(s, steps, y=Inches(3.0), h=Inches(1.1), color=NAVY, size=12.5):
    n = len(steps)
    gap = Inches(0.12)
    total = Inches(12.13)
    w = Emu(int((total - gap * (n - 1)) / n))
    x = Inches(0.6)
    for i, st in enumerate(steps):
        shp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
        shp.adjustments[0] = 0.35
        shp.fill.solid()
        shp.fill.fore_color.rgb = color if i < n - 1 else ORANGE
        shp.line.fill.background()
        shp.shadow.inherit = False
        shape_text(shp, [(st, size, True, WHITE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(int(x + w + gap))


def kpi_cards(s, items, y=Inches(2.4), h=Inches(2.3)):
    n = len(items)
    gap = Inches(0.3)
    total = Inches(12.13)
    w = Emu(int((total - gap * (n - 1)) / n))
    x = Inches(0.6)
    for label, big, sub, color in items:
        card = rect(s, x, y, w, h, LIGHT)
        shape_text(card, [
            (label, 12.5, True, GRAY, PP_ALIGN.CENTER),
            (big, 34, True, color, PP_ALIGN.CENTER),
            (sub, 11.5, False, INK, PP_ALIGN.CENTER),
        ], anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(int(x + w + gap))


def two_panel(s, left, right, y=Inches(2.0), h=Inches(3.9)):
    for (title, col, lines), x in ((left, Inches(0.6)), (right, Inches(6.85))):
        card = rect(s, x, y, Inches(5.9), h, WHITE, line=col)
        body = [(title, 16, True, col, PP_ALIGN.LEFT)]
        for ln in lines:
            body.append(("• " + ln, 12.5, False, INK, PP_ALIGN.LEFT))
        shape_text(card, body)


def rows_table(s, rows, y=Inches(2.0), rh=Inches(0.86), highlight_idx=None):
    highlight_idx = highlight_idx or set()
    yy = y
    for i, (a, b) in enumerate(rows):
        hot = i in highlight_idx
        card = rect(s, Inches(0.6), yy, Inches(12.13), rh, SOFT if hot else LIGHT, line=ORANGE if hot else None)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.18)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = a + "   "
        f = r1.font; f.name, f.size, f.bold = FONT, Pt(14), True
        f.color.rgb = ORANGED if hot else NAVY; _ea(r1)
        r2 = p.add_run(); r2.text = b
        f = r2.font; f.name, f.size, f.bold = FONT, Pt(12.5), False; f.color.rgb = INK; _ea(r2)
        yy = Emu(int(yy + rh + Inches(0.12)))


def section_divider(text, sub=""):
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY, round_=False)
    rect(s, Inches(0.9), Inches(3.0), Inches(0.16), Inches(1.4), ORANGE, round_=False)
    put_text(s, Inches(1.25), Inches(3.0), Inches(11), Inches(1.5),
             [(text, 34, True, WHITE, PP_ALIGN.LEFT),
              (sub, 16, False, RGBColor(0xC7, 0xD2, 0xE4), PP_ALIGN.LEFT)])
    return s


# ════════════════════════════════════════════════════════════
# S1. 표지
# ════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY, round_=False)
rect(s, Inches(0.9), Inches(2.4), Inches(0.16), Inches(1.7), ORANGE, round_=False)
put_text(s, Inches(1.3), Inches(2.0), Inches(11), Inches(0.5),
         [("고객사 공유회 · Proof of Concept", 15, True, ORANGE, PP_ALIGN.LEFT)])
put_text(s, Inches(1.3), Inches(2.55), Inches(11.6), Inches(1.7),
         [("한화투자증권 AI 상담 어시스턴트", 42, True, WHITE, PP_ALIGN.LEFT),
          ("상담사를 돕고(응대) · 가르치고(코치) · 스스로 좋아지는(제보) AI", 18, False, RGBColor(0xC7, 0xD2, 0xE4), PP_ALIGN.LEFT)])
put_text(s, Inches(1.3), Inches(5.8), Inches(11), Inches(0.9),
         [("업무 편람을 근거로 답하고, 출처까지 보여주는 상담 도우미", 14, True, WHITE, PP_ALIGN.LEFT),
          ("증권 ITO팀  ·  발표 약 30분", 12.5, False, RGBColor(0x9C, 0xA8, 0xBC), PP_ALIGN.LEFT)])

# ════════════════════════════════════════════════════════════
# S2. 아젠다 (쉬운 말)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "오늘 이야기", "30분 동안 이 순서로 말씀드립니다")
items = [
    ("01", "왜 만들었나", "상담 현장의 두 가지 문제 — 검색에 드는 시간, 신입 교육"),
    ("02", "무엇을 만들었나", "한눈에 보는 전체 그림 + 3가지 기능(응대·AI코치·제보)"),
    ("03", "어떻게 작동하나", "쉬운 비유로 원리 + (전문) 아키텍처·검색 기법"),
    ("04", "얼마나 정확한가", "정확도 수치와 '어떻게 쟀는지' 근거"),
    ("05", "한계와 앞으로", "AI의 한계·해결 방향 + 데이터/음성 비전 + 도입"),
]
y = Inches(1.85)
for num, t, d in items:
    card = rect(s, Inches(0.6), y, Inches(12.13), Inches(0.92), LIGHT)
    shape_text(card, [(f"{num}   {t}", 17, True, NAVY, PP_ALIGN.LEFT),
                      (d, 12, False, GRAY, PP_ALIGN.LEFT)], anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(1.04)))

# ════════════════════════════════════════════════════════════
# S3. ★ 한눈에 보는 전체 그림 (big picture)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "한눈에 보기", "전체 그림 — 이 한 장만 기억하셔도 됩니다")
# 좌: 상담사 3가지 사용  / 중앙: AI 두뇌 / 우: 지식 + 폐쇄망
c1 = rect(s, Inches(0.6), Inches(1.8), Inches(3.7), Inches(3.9), WHITE, line=NAVY)
shape_text(c1, [("상담사가 이렇게 씁니다", 14, True, NAVY, PP_ALIGN.LEFT),
                ("", 5, False, INK, PP_ALIGN.LEFT),
                ("💬 응대 모드", 14, True, ORANGE, PP_ALIGN.LEFT),
                ("  고객 응대 중 즉시 검색·답변", 11.5, False, INK, PP_ALIGN.LEFT),
                ("🎓 AI 코치", 14, True, ORANGE, PP_ALIGN.LEFT),
                ("  다양한 고객 상황을 연습·채점", 11.5, False, INK, PP_ALIGN.LEFT),
                ("⚑ 오답 제보", 14, True, ORANGE, PP_ALIGN.LEFT),
                ("  틀린 답 신고 → 품질 개선", 11.5, False, INK, PP_ALIGN.LEFT)])
# 화살표
a1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.45), Inches(3.4), Inches(0.7), Inches(0.6))
a1.fill.solid(); a1.fill.fore_color.rgb = ORANGE; a1.line.fill.background(); a1.shadow.inherit = False
c2 = rect(s, Inches(5.3), Inches(1.8), Inches(2.7), Inches(3.9), SOFT, line=ORANGE)
shape_text(c2, [("AI 두뇌", 14, True, ORANGED, PP_ALIGN.CENTER),
                ("", 6, False, INK, PP_ALIGN.CENTER),
                ("① 질문의 '의미'로", 12, False, INK, PP_ALIGN.CENTER),
                ("편람에서 근거 검색", 12, True, NAVY, PP_ALIGN.CENTER),
                ("", 4, False, INK, PP_ALIGN.CENTER),
                ("② 찾은 근거로", 12, False, INK, PP_ALIGN.CENTER),
                ("답변 생성 + 출처", 12, True, NAVY, PP_ALIGN.CENTER),
                ("", 6, False, INK, PP_ALIGN.CENTER),
                ("(RAG 방식)", 11, False, GRAY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
a2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.15), Inches(3.4), Inches(0.7), Inches(0.6))
a2.fill.solid(); a2.fill.fore_color.rgb = ORANGE; a2.line.fill.background(); a2.shadow.inherit = False
c3 = rect(s, Inches(9.0), Inches(1.8), Inches(3.73), Inches(3.9), WHITE, line=TEAL)
shape_text(c3, [("지식 + 보안", 14, True, TEAL, PP_ALIGN.LEFT),
                ("", 5, False, INK, PP_ALIGN.LEFT),
                ("📚 업무 편람(위키)", 13, True, NAVY, PP_ALIGN.LEFT),
                ("  매일 자동 최신화", 11.5, False, INK, PP_ALIGN.LEFT),
                ("🔎 의미 검색 색인", 13, True, NAVY, PP_ALIGN.LEFT),
                ("  회사 안에서 무료로 구동", 11.5, False, INK, PP_ALIGN.LEFT),
                ("🔒 폐쇄망 대비 설계", 13, True, NAVY, PP_ALIGN.LEFT),
                ("  인터넷 없이도 동작", 11.5, False, INK, PP_ALIGN.LEFT)])
terms(s, [("RAG", "AI가 편람을 '찾아 읽고' 답하는 방식(외워서가 아니라 오픈북 시험처럼)"),
          ("폐쇄망", "보안을 위해 인터넷과 분리해 둔 사내 전용망(망분리)")])

# ════════════════════════════════════════════════════════════
# S4. 문제정의
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "01 · 왜 만들었나", "현장의 문제 — 검색에 쓰는 시간, 맨몸으로 배우는 신입")
kpi_cards(s, [
    ("고객 1건 응대 중 편람 검색", "120~300초", "숙련자도 직접 찾음", NAVY),
    ("신입이 혼자 응대까지", "약 3개월", "선배가 1:1로 붙어야", NAVY),
    ("지식이 흩어진 곳", "여기저기", "편람·FAQ·선배 머릿속", ORANGE),
], y=Inches(2.2), h=Inches(2.5))
put_text(s, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.1),
         [("• 정답은 편람·FAQ·선배 머릿속에 흩어져 있어, 찾는 일 자체가 시간을 잡아먹습니다.", 14, False, INK, PP_ALIGN.LEFT),
          ("• 콜센터는 근무시간의 상당 부분을 '검색'에 씁니다 (McKinsey·IDC: 정보탐색에 근무시간 20~30%).", 13, False, GRAY, PP_ALIGN.LEFT)])
takeaway(s, "검색도 교육도 '사람의 시간'을 쓰는 구조 → AI로 응대 지원·훈련·품질개선을 함께")

# ════════════════════════════════════════════════════════════
# S5. 솔루션 3기능
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "02 · 무엇을 만들었나", "하나의 편람으로, 세 가지 쓰임")
cards = [
    ("💬 응대 모드", ORANGE, ["고객 응대 중 편람을 즉시 검색",
                          "질문 → 출처·신뢰도 먼저, 답변은 이어서",
                          "출처를 누르면 편람(위키)으로 바로 이동"]),
    ("🎓 AI 코치", NAVY, ["다양한 고객 상황을 가정해 응대 연습",
                        "커리큘럼·복습으로 체계적 학습",
                        "응대를 점수·코칭으로 피드백"]),
    ("⚑ 오답 제보", TEAL, ["답이 틀리면 상담사가 바로 신고",
                        "사유·정답 제안과 함께 축적",
                        "검토 → 편람/답변 개선 (쓸수록 좋아짐)"]),
]
n = len(cards); gap = Inches(0.3); w = Emu(int((Inches(12.13) - gap * (n - 1)) / n)); x = Inches(0.6)
for title, col, lines in cards:
    card = rect(s, x, Inches(2.0), w, Inches(3.7), WHITE, line=col)
    body = [(title, 17, True, col, PP_ALIGN.LEFT)]
    for ln in lines:
        body.append(("• " + ln, 12.5, False, INK, PP_ALIGN.LEFT))
    shape_text(card, body)
    x = Emu(int(x + w + gap))
takeaway(s, "같은 편람을 — 응대는 '검색→답변', 코치는 '연습→채점', 제보는 '신고→개선'으로")

# ════════════════════════════════════════════════════════════
# S6. 데모 ① 응대 모드
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "02 · 무엇을 만들었나", "데모 ① 응대 모드 — 근거와 함께 답한다")
two_panel(s,
          ("화면에서 일어나는 일", NAVY, ["상담사가 평소 말로 질문",
                                  "검색 끝나면 '출처+신뢰도'가 먼저 뜸",
                                  "답변이 글자 단위로 이어서 표시",
                                  "출처 카드를 누르면 위키 페이지로 이동",
                                  "답이 이상하면 '오답 제보' 한 번"]),
          ("상담사에게 좋은 점", ORANGE, ["근거(출처)를 보고 바로 신뢰·확인",
                                   "편람에 없으면 '없다'고 정직하게 답",
                                   "신뢰도(높음/보통/낮음)로 주의 환기",
                                   "추후 음성(STT)으로 타이핑도 줄임"]))
terms(s, [("출처", "답변의 근거가 된 편람 페이지 — 눌러서 원문 확인 가능"),
          ("신뢰도", "검색이 얼마나 확실한지 색으로 표시(높음/보통/낮음)")])

# ════════════════════════════════════════════════════════════
# S7. 데모 ② AI 코치 (최신: 뱅크/커리큘럼/복습/페르소나/음성)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "02 · 무엇을 만들었나", "데모 ② AI 코치 — '어떤 고객이 올지 모른다'를 연습")
two_panel(s,
          ("체계적 학습 구조", NAVY, ["고객 상황 8종(화남·급함·고령·불안 등)",
                              "난이도(초·중·고)별 시나리오",
                              "커리큘럼: 순서대로 / 랜덤: 실전 감각",
                              "복습: 틀린 것만 모아 다시(오답 노트)",
                              "고객 질문을 '음성'으로 읽어줌(TTS)"]),
          ("응대 후 받는 피드백", ORANGE, ["종합 점수 + 항목별 막대",
                                   "포함/누락한 핵심 안내 표시",
                                   "'다음엔 이렇게' 실행 코칭",
                                   "모범 답변 + 편람 출처(링크)"]))
terms(s, [("페르소나", "고객 유형(화난·급한·고령 고객 등)을 골라 상황을 가정"),
          ("TTS", "글로 된 시나리오를 사람 목소리로 읽어주는 기술")])

# ════════════════════════════════════════════════════════════
# S8. 데모 ③ 오답 제보 → 검토 (선순환)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "02 · 무엇을 만들었나", "데모 ③ 오답 제보 → 검토 — 쓸수록 좋아진다")
flow(s, ["답변 의심", "한 번에 제보\n(사유+정답제안)", "자동 축적", "운영자 검토", "편람·답변 개선"],
     y=Inches(2.2), h=Inches(1.15), size=12)
put_text(s, Inches(0.9), Inches(3.9), Inches(11.6), Inches(1.8),
         [("• 제보 시 질문·답변·출처를 '그 순간 그대로' 저장 → 나중에 재현·추적 가능", 14, False, INK, PP_ALIGN.LEFT),
          ("• 검토 화면에서 미처리/처리완료로 관리, 외부 인터넷 없이 사내 파일로 보관", 14, False, INK, PP_ALIGN.LEFT),
          ("• 이 '신고→개선' 고리가 곧 데이터가 쌓일수록 똑똑해지는 선순환의 출발점", 14, True, ORANGED, PP_ALIGN.LEFT)])
takeaway(s, "틀린 답을 '버그'가 아니라 '개선 데이터'로 바꾼다", color=TEAL)

# ════════════════════════════════════════════════════════════
# 섹션 구분 — 작동 원리
# ════════════════════════════════════════════════════════════
section_divider("03 · 어떻게 작동하나", "먼저 쉬운 비유로, 그다음 전문적으로")

# ════════════════════════════════════════════════════════════
# S10. 쉬운 비유 (non-IT)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "03 · 어떻게 작동", "쉽게 말하면 — '오픈북 시험을 보는 AI'")
flow(s, ["고객 질문", "질문의 '의미'로\n편람에서 검색", "관련 편람 몇 장\n찾아오기", "그것만 읽고\n답 작성", "답변 + 출처"],
     y=Inches(2.2), h=Inches(1.2), size=11.5)
c = rect(s, Inches(0.6), Inches(3.8), Inches(12.13), Inches(2.0), LIGHT)
shape_text(c, [("핵심: AI가 '외워서' 답하지 않습니다. 매번 편람을 '찾아 읽고' 답합니다.", 15, True, NAVY, PP_ALIGN.LEFT),
               ("• 편람이 바뀌면? → AI를 다시 가르칠 필요 없이, 새 편람만 넣으면 끝", 13.5, False, INK, PP_ALIGN.LEFT),
               ("• '어디서 나온 답?' → 찾아 읽은 편람이 곧 출처 (금융에선 필수)", 13.5, False, INK, PP_ALIGN.LEFT),
               ("• 편람에 없으면 '없다'고 답하게 만들어 → 그럴듯한 거짓말(환각) 억제", 13.5, False, INK, PP_ALIGN.LEFT)])
terms(s, [("환각", "AI가 사실이 아닌데도 그럴듯하게 지어내는 현상")])

# ════════════════════════════════════════════════════════════
# S11. 아키텍처 (전문 청중)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "03 · 어떻게 작동", "(전문) 시스템 구성 — 폐쇄망 이식을 전제로")
flow(s, ["상담사 화면\n(React)", "FastAPI\n(서버)", "의미 검색\n(ChromaDB)", "답변 생성\n(LLM)", "출처+답변"],
     y=Inches(2.0), h=Inches(1.15), size=11.5)
rows = [
    ("프론트  React", "응대/AI코치/제보검토 3화면, 출처·음성·마크다운 표시"),
    ("백엔드  FastAPI", "검색·채점·제보 로직 분리, 모델 호출 창구 통일(교체 대비)"),
    ("검색  ChromaDB", "의미검색 색인. 파일 기반 → 폴더 복사만으로 이식"),
    ("임베딩  bge-m3", "회사 안에서 무료·오프라인 구동(의미를 숫자로). 생성 LLM만 교체 지점"),
]
rows_table(s, rows, y=Inches(3.5), rh=Inches(0.66), highlight_idx={3})
terms(s, [("임베딩", "글의 '의미'를 숫자로 바꿔 비슷한 내용을 찾게 하는 기술"),
          ("LLM", "사람처럼 글을 이해·생성하는 거대 AI(예: Gemini, 사내모델)")])

# ════════════════════════════════════════════════════════════
# S12. 차별화 검색 기법 요약 (쉬운 설명 + 각주)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "03 · 어떻게 작동", "정확도를 높인 검색 기법 — 일반적 방식과 다른 점")
rows_table(s, [
    ("제목·본문 따로 보기", "제목과 본문을 따로 비교 → 짧은 질문(예: '양도세')에도 정확"),
    ("같은 문서 중복 정리", "여러 표현을 달아도 한 문서가 부당하게 유리해지지 않게(Max Pooling)"),
    ("의미+단어 동시 검색", "의미검색에 키워드검색을 더함 → 'K-OTC','OTP' 같은 고유명사 강함"),
    ("동의어 자동 보강", "'미국주식→해외주식' 등 사내 용어를 미리 연결(오프라인 사전)"),
    ("신뢰도 자동 보정", "실측 점수 분포로 '높음/보통/낮음' 기준을 데이터로 조정"),
], y=Inches(1.9), rh=Inches(0.78), highlight_idx={0, 1, 2, 3, 4})
note(s, "* 일반 챗봇은 보통 한 가지 방식만 씁니다. 우리는 5가지를 조합해 검색 정확도를 끌어올렸습니다.")

# ════════════════════════════════════════════════════════════
# 섹션 — 정확도
# ════════════════════════════════════════════════════════════
section_divider("04 · 얼마나 정확한가", "수치보다 '어떻게 쟀는지'가 중요합니다")

# ════════════════════════════════════════════════════════════
# S14. 측정 방법론
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "04 · 얼마나 정확한가", "어떻게 쟀나 — 검색과 답변을 '나눠서' 측정")
two_panel(s,
          ("① 검색 정확도", NAVY, ["질문의 정답 문서가 검색 상위에 들어오나?",
                              "지표: 정답이 상위 3개 안에 든 비율(hit@3)",
                              "30문항 평가셋으로 측정",
                              "→ 측정 결과 파일로 '재현 가능'"]),
          ("② 답변 정확도", ORANGE, ["생성된 답변이 정답 내용을 담았나?",
                              "사람 대신 AI가 '의미로' 채점(표현 달라도 OK)",
                              "→ '단어만 같은지'가 아니라 '뜻이 맞는지'",
                              "측정 방법이 정해져 있어 반복 측정 가능"]))
terms(s, [("hit@3", "정답 문서가 검색 결과 상위 3개 안에 든 비율(높을수록 좋음)"),
          ("평가셋", "정답을 미리 정해 둔 시험 문제 묶음(여기선 30문항)")])

# ════════════════════════════════════════════════════════════
# S15. 검색 정확도 (증빙)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "04 · 얼마나 정확한가", "검색 정확도 — 정답을 거의 항상 찾아온다")
kpi_cards(s, [
    ("hit@1 (1순위 정답)", "70%", "정답이 맨 위", NAVY),
    ("hit@3 (상위 3개 안)", "100%", "정답이 top-3 안", ORANGE),
    ("hit@5 (상위 5개 안)", "100%", "정답이 top-5 안", NAVY),
    ("MRR (평균 순위 점수)", "0.84", "1에 가까울수록 좋음", NAVY),
], y=Inches(2.2), h=Inches(2.3))
put_text(s, Inches(0.9), Inches(4.8), Inches(11.6), Inches(1.1),
         [("• 30문항 기준, 정답 편람이 '항상' 검색 상위 3개 안에 들어옵니다(hit@3 100%).", 14, False, INK, PP_ALIGN.LEFT),
          ("• 측정 결과는 파일로 저장되어 그대로 재현·검증할 수 있습니다.", 13, False, GRAY, PP_ALIGN.LEFT)])
note(s, "근거: data/eval/accuracy_bge_m3_search_*.json (30문항, 회사 안 로컬 검색 기준)")

# ════════════════════════════════════════════════════════════
# S16. 답변 정확도 (방법론 + 정직)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "04 · 얼마나 정확한가", "답변 정확도 — 숫자보다 '방법'을 정직하게")
two_panel(s,
          ("측정값", ORANGE, ["AI 의미 채점: 12/15 = 80%",
                          "단어만 보는 채점(하한)은 더 낮게 나옴",
                          "→ 표현이 다른 정답을 놓치기 때문",
                          "→ 그래서 '의미 채점'이 더 정확"]),
          ("정직한 한계", NAVY, ["답변 정확도 100%는 불가능(AI 특성)",
                            "확실치 않으면 '확인되지 않음'으로 답",
                            "고위험 답변은 사람(상담원)이 최종 확인",
                            "측정 방법이 문서화되어 언제든 재측정"]))
takeaway(s, "검색은 '하드 증빙', 답변은 '방법론+재현'으로 — 과장 없이 보여드립니다")

# ════════════════════════════════════════════════════════════
# 섹션 — 한계와 앞으로
# ════════════════════════════════════════════════════════════
section_divider("05 · 한계와 앞으로", "AI의 한계, 업계의 해결, 그리고 우리의 비전")

# ════════════════════════════════════════════════════════════
# S18. LLM 한계 ① 환각
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "한계 ① 환각 — '그럴듯하게 틀린 답'")
two_panel(s,
          ("업계의 방어 4단계", NAVY, ["① 편람 안에서만 답하게(그라운딩)",
                                "② 답에 출처를 꼭 붙이기",
                                "③ 생성 후 근거와 맞는지 재확인",
                                "④ 자신 없으면 '모른다'로 거부",
                                "'오답은 모른다고 하는 것보다 나쁘다'"]),
          ("우리 적용", ORANGE, ["편람 근거 + 출처 + 신뢰도 + 거부",
                            "오답 제보로 사람이 교정(휴먼인더루프)",
                            "향후: 답변 근거 일치 자동 점검(RAGAS)",
                            "현실: RAG도 환각을 70~90%만 줄임 → 사람 검토 병행"]))
note(s, "근거: OpenAI 'why LMs hallucinate', 스탠퍼드 법률RAG(잔존 17~33%) — RESEARCH_NOTES 참조")

# ════════════════════════════════════════════════════════════
# S19. LLM 한계 ② 최신성 / ③ 도메인
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "한계 ② 최신성 · ③ 전문성")
two_panel(s,
          ("② 최신성(옛 정보)", NAVY, ["AI는 학습 시점 이후를 모름",
                               "해결: 다시 가르치지 않고 '편람을 갱신'",
                               "변경분만 매일 자동 반영(증분)",
                               "시세·공시는 도구 연결로 보완(향후)"]),
          ("③ 전문성(증권 도메인)", ORANGE, ["범용 AI는 증권 약관·수치에 약함",
                                   "RAG(편람 검색)로 정확도 확보가 1순위",
                                   "검색 설계가 모델보다 중요(연구: 56% vs 19%)",
                                   "향후: 응대 말투·포맷만 가볍게 미세조정"]))
note(s, "근거: Agentic RAG(arXiv 2501.09136), FinanceBench, Red Hat RAG vs Fine-tuning — RESEARCH_NOTES 참조")

# ════════════════════════════════════════════════════════════
# S20. LLM 한계 ④ 보안/망분리
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "한계 ④ 보안·망분리 — 한국 금융 특유의 변수")
rows_table(s, [
    ("규제 환경", "금융위 '망분리 개선 로드맵'(2024-08) — 74개사 141건 혁신금융 특례 신청"),
    ("국산 모델 선택지", "HyperCLOVA X(미래에셋 온프레미스)·EXAONE·Qwen3(카카오뱅크)·Solar(신한DS)"),
    ("핵심 설계", "처음부터 폐쇄망 이식형 — 인터넷 없이 사내에서 구동 가능"),
    ("거버넌스", "사람 검토(휴먼인더루프) + 출처 추적성 + 감사 로그"),
], y=Inches(2.0), rh=Inches(0.92), highlight_idx={2})
takeaway(s, "우리는 폐쇄망을 보고 설계 — 도입 시 '바꿀 곳은 답변 AI 단 1곳'", color=TEAL)

# ════════════════════════════════════════════════════════════
# S21. 비전 ① 데이터 허브
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "비전 ① 데이터 허브 — 흩어진 데이터를 답변으로")
put_text(s, Inches(0.9), Inches(1.55), Inches(11.6), Inches(0.5),
         [("지금: 업무편람 한 가지  →  앞으로: 메일·메신저·엑셀·회의록까지 모아 답변 DB로", 14, True, NAVY, PP_ALIGN.LEFT)])
rows_table(s, [
    ("일본 SMBC", "사내 규정·매뉴얼 130만 건을 AI로 통합 검색"),
    ("일본 미즈호 / 일본생명", "사무수속 RAG 해결률 96% · 약관 RAG 유효답변 90%(엑셀 표까지 판독)"),
    ("국내 신한투자·미래에셋·키움", "사내문서 RAG · 온프레미스 사내 AI 챗봇 운영"),
    ("기술 흐름", "수집 → 정리 → 핵심표시 → 중복제거 → 개인정보 가림 → 권한 보존"),
], y=Inches(2.15), rh=Inches(0.86), highlight_idx={3})
note(s, "출처: RESEARCH_NOTES_2026-06-20.md (각 사 IR/보도). 비IT 요약: '회사의 흩어진 지식을 한곳에서 답하게'")

# ════════════════════════════════════════════════════════════
# S22. 비전 ② 음성(STT/TTS)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "비전 ② 음성 연동 — 듣고(STT) 말하는(TTS) 상담")
two_panel(s,
          ("STT — 통화를 글로", NAVY, ["고객 통화를 실시간으로 받아쓰기",
                               "상담사가 타이핑 없이 자동 질의",
                               "응대에 더 집중, 통화-답변 기록도 축적"]),
          ("TTS — 시나리오를 음성으로", ORANGE, ["AI 코치의 고객 질문을 목소리로 재생",
                                      "고객 유형별 말투(화남·고령 등) 반영",
                                      "엔진 교체형: 사내 음성엔진으로 대체 가능"]))
terms(s, [("STT", "사람 말을 글자로 바꾸는 기술(받아쓰기)"),
          ("TTS", "글자를 사람 목소리로 바꾸는 기술(읽어주기)")])

# ════════════════════════════════════════════════════════════
# S23. 폐쇄망 로드맵
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "실도입(폐쇄망) — 바꿀 곳은 단 한 곳")
rows_table(s, [
    ("답변 생성 AI", "외부 API → 사내 모델(예: Qwen3 / HyperCLOVA X / Solar)  ★여기만 교체"),
    ("의미 검색 색인", "회사 안 로컬(bge-m3) — 그대로 (무료·오프라인)"),
    ("검색 저장소", "파일 기반 ChromaDB — 폴더 복사로 이식 끝"),
    ("화면·서버", "React + FastAPI — 변경 없음"),
], y=Inches(2.0), rh=Inches(0.92), highlight_idx={0})
note(s, "편람은 매일 자동 최신화(변경분만 반영) → 운영 부담 최소, 외부 인터넷 의존 0")

# ════════════════════════════════════════════════════════════
# S24. ROI
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "기대 효과 — 검색 시간 절감")
kpi_cards(s, [
    ("편람 검색 체감", "300초 → 3초", "응답 속도", ORANGE),
    ("분기 절감(추정)", "2.7~3.8 M/M", "검색시간 35%↓ 가정", NAVY),
    ("정성 효과", "균일한 품질", "신입 온보딩↓·지식 자산화", NAVY),
], y=Inches(2.3), h=Inches(2.3))
put_text(s, Inches(0.9), Inches(4.9), Inches(11.6), Inches(1.0),
         [("산정 예시: 39명 × 검색 1.8~2.5h/일 × 53영업일 × 35%↓ ≈ 분기 2.7~3.8 M/M", 13.5, False, INK, PP_ALIGN.LEFT),
          ("* 인력·콜 수는 예시 파라미터. 실도입 시 한화 실측치로 재산정.", 12, False, GRAY, PP_ALIGN.LEFT)])
takeaway(s, "M/M = Man-Month(한 사람이 한 달 일하는 양) — 분기당 약 3명 분의 시간 절약")

# ════════════════════════════════════════════════════════════
# S25. 도입 로드맵
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · 한계와 앞으로", "제안 — 단계별 도입")
flow(s, ["① PoC 검증", "② 응대 모드 확대", "③ 코치·제보 운영", "④ 데이터허브·음성", "⑤ 폐쇄망 본도입"],
     y=Inches(2.4), h=Inches(1.15), size=12)
put_text(s, Inches(0.9), Inches(4.1), Inches(11.6), Inches(1.6),
         [("• 작게 검증(한화 편람 일부) → 효과를 숫자로 확인", 14, False, INK, PP_ALIGN.LEFT),
          ("• 응대 실사용 + 오답 제보로 답변 품질을 키움", 14, False, INK, PP_ALIGN.LEFT),
          ("• 비정형 데이터·음성 연동, 사내 모델로 폐쇄망 안착", 14, False, INK, PP_ALIGN.LEFT)])
takeaway(s, "작게 검증하고, 데이터로 키우고, 폐쇄망으로 안착")

# ════════════════════════════════════════════════════════════
# S26. 요청사항
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "마무리", "한화 측에 요청드리는 것")
rows_table(s, [
    ("데이터 샘플", "업무편람·FAQ 일부(정확도 실측용) + (가능 시) 상담 로그 비식별 샘플"),
    ("음성 연동", "콜 인프라/녹취 연동 가능 범위 협의"),
    ("모델·인프라", "폐쇄망 GPU·사내 모델 후보(HyperCLOVA X/Qwen/Solar)·망분리 정책 확인"),
], y=Inches(2.1), rh=Inches(0.95))
takeaway(s, "작은 데이터로 빠르게 실측 → 효과를 숫자로 확인하고 다음 단계 결정")

# ════════════════════════════════════════════════════════════
# S27. 용어 사전 (비IT)
# ════════════════════════════════════════════════════════════
s = slide()
header(s, "참고", "용어 사전 — 쉬운 한 줄 설명")
gloss = [
    ("RAG", "AI가 편람을 찾아 읽고 답하는 방식(오픈북 시험처럼)"),
    ("LLM", "사람처럼 글을 이해·생성하는 거대 AI(예: Gemini)"),
    ("임베딩", "글의 '의미'를 숫자로 바꿔 비슷한 내용을 찾게 함"),
    ("벡터DB(ChromaDB)", "의미로 검색하는 색인 저장소(파일로 이식 가능)"),
    ("bge-m3", "회사 안에서 무료로 돌리는 의미검색용 AI"),
    ("폐쇄망(망분리)", "보안 위해 인터넷과 분리한 사내 전용망"),
    ("환각", "AI가 사실이 아닌데 그럴듯하게 지어내는 현상"),
    ("출처/신뢰도", "답의 근거 편람 / 검색이 얼마나 확실한지 표시"),
    ("hit@3 / MRR", "정답이 상위 3개 안에 든 비율 / 평균 순위 점수"),
    ("STT / TTS", "말→글 받아쓰기 / 글→말 읽어주기"),
]
y = Inches(1.75)
for i, (t, d) in enumerate(gloss):
    col = 0 if i < 5 else 1
    row = i % 5
    x = Inches(0.6) if col == 0 else Inches(6.85)
    yy = Emu(int(Inches(1.75) + Inches(0.98) * row))
    card = rect(s, x, yy, Inches(5.9), Inches(0.86), LIGHT)
    shape_text(card, [(t, 13.5, True, NAVY, PP_ALIGN.LEFT),
                      (d, 11.5, False, INK, PP_ALIGN.LEFT)], anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════
# S28. 클로징
# ════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY, round_=False)
put_text(s, Inches(1.25), Inches(1.7), Inches(11), Inches(3.6),
         [("세 줄 요약", 18, True, ORANGE, PP_ALIGN.LEFT),
          ("1.  상담을 돕고(응대)·가르치고(코치)·스스로 좋아지는(제보) AI — 하나의 편람으로", 17, True, WHITE, PP_ALIGN.LEFT),
          ("2.  정확도는 구조에서 나온다 — 검색 hit@3 100%, 출처로 신뢰 확보", 17, True, WHITE, PP_ALIGN.LEFT),
          ("3.  처음부터 폐쇄망 이식형 — 도입 시 바꿀 곳은 답변 AI 단 1곳", 17, True, WHITE, PP_ALIGN.LEFT)])
put_text(s, Inches(1.25), Inches(5.9), Inches(11), Inches(0.8),
         [("Q&A   ·   상세 근거: docs/api-spec.md · ANSWER_QUALITY.md · ONPREM_ROADMAP.md · RESEARCH_NOTES_2026-06-20.md", 12.5, False, RGBColor(0xC7, 0xD2, 0xE4), PP_ALIGN.LEFT)])

# ── 저장 ────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "한화투자증권_AI상담_PoC_30min.pptx")
prs.save(out)
print(f"저장 완료: {out}  (슬라이드 {len(prs.slides._sldIdLst)}장)")
